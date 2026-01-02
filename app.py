import streamlit as st
import pandas as pd
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
from datetime import datetime

# Configure page
st.set_page_config(
    page_title="AI Excel & PowerPoint Agent",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS
st.markdown("""
    <style>
    .main-title {
        font-size: 2.5rem;
        color: #1f77b4;
        margin-bottom: 1rem;
    }
    </style>
""", unsafe_allow_html=True)

# Sidebar
st.sidebar.title("⚙️ Configuration")
st.sidebar.markdown("---")

mode = st.sidebar.radio(
    "Select Mode:",
    ["📈 Excel Automation", "📊 PowerPoint Creation", "🔄 Excel → PowerPoint"]
)

# Helper functions
def create_ppt_from_data(data_dict: dict, title: str, theme_color: tuple) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme_color)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Content slides
    for section_title, content in data_dict.items():
        if section_title == "title":
            continue
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = section_title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        
        body = slide.placeholders[1].text_frame
        body.clear()
        
        if isinstance(content, list):
            for item in content:
                p = body.add_paragraph()
                p.text = str(item)[:100]
                p.font.size = Pt(18)
        elif isinstance(content, dict):
            for k, v in content.items():
                p = body.add_paragraph()
                p.text = f"{k}: {v}"
                p.font.size = Pt(16)
    
    return prs

# Main UI
st.markdown("<h1 class='main-title'>📊 AI Excel & PowerPoint Agent</h1>", unsafe_allow_html=True)
st.markdown("Powered by Streamlit")
st.markdown("---")

if mode == "📈 Excel Automation":
    st.header("📈 Excel Data Analysis")
    uploaded_file = st.file_uploader("Upload Excel or CSV file", type=["xlsx", "xls", "csv"])
    
    if uploaded_file:
        try:
            df = pd.read_csv(uploaded_file) if uploaded_file.name.endswith('.csv') else pd.read_excel(uploaded_file)
            st.success(f"✅ Loaded: {df.shape[0]} rows, {df.shape[1]} columns")
            
            # Display preview
            st.subheader("Data Preview")
            st.dataframe(df.head(10), use_container_width=True)
            
            # Basic statistics
            st.subheader("Data Statistics")
            col1, col2, col3 = st.columns(3)
            col1.metric("Total Rows", df.shape[0])
            col2.metric("Total Columns", df.shape[1])
            col3.metric("Memory Usage (KB)", f"{df.memory_usage(deep=True).sum() / 1024:.2f}")
            
            # Data description
            st.subheader("Data Description")
            st.write(df.describe())
            
            # Download processed file
            output_buffer = io.BytesIO()
            with pd.ExcelWriter(output_buffer, engine='openpyxl') as writer:
                df.to_excel(writer, sheet_name='Data', index=False)
            output_buffer.seek(0)
            
            st.download_button(
                label="💾 Download Processed Excel",
                data=output_buffer,
                file_name=f"processed_{uploaded_file.name}",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )
        except Exception as e:
            st.error(f"Error: {str(e)}")

elif mode == "📊 PowerPoint Creation":
    st.header("📊 Generate PowerPoint Presentation")
    
    col1, col2 = st.columns([2, 1])
    with col1:
        uploaded_file = st.file_uploader("Upload Excel for PPT", type=["xlsx", "xls", "csv"])
    with col2:
        theme_color_name = st.selectbox("Theme Color", ["Blue", "Green", "Red", "Purple", "Orange"])
    
    colors = {"Blue": (31, 119, 180), "Green": (44, 160, 44), "Red": (214, 39, 40), "Purple": (148, 103, 189), "Orange": (255, 127, 14)}
    ppt_title = st.text_input("Presentation Title", "Data Report")
    
    if uploaded_file and st.button("🎨 Generate PowerPoint"):
        try:
            df = pd.read_csv(uploaded_file) if uploaded_file.name.endswith('.csv') else pd.read_excel(uploaded_file)
            st.info(f"Creating presentation from {df.shape[0]} rows of data...")
            
            ppt_content = {
                "Data Overview": [
                    f"Total Records: {df.shape[0]}",
                    f"Total Fields: {df.shape[1]}",
                    f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
                ],
                "Key Columns": list(df.columns[:6]),
                "Data Summary": {
                    "Rows": df.shape[0],
                    "Columns": df.shape[1],
                    "Data Types": len(df.dtypes.unique())
                }
            }
            
            prs = create_ppt_from_data(ppt_content, ppt_title, colors[theme_color_name])
            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            st.success("✅ PowerPoint generated successfully!")
            st.download_button(
                label="📥 Download PowerPoint",
                data=ppt_buffer,
                file_name=f"{ppt_title.replace(' ', '_')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except Exception as e:
            st.error(f"Error: {str(e)}")

else:  # End-to-End Workflow
    st.header("🔄 Complete Excel to PowerPoint Workflow")
    
    uploaded_file = st.file_uploader("Upload Excel file", type=["xlsx", "xls", "csv"])
    ppt_title = st.text_input("Presentation Title", "Analysis Report")
    
    if uploaded_file and st.button("🚀 Create Complete Presentation"):
        try:
            df = pd.read_csv(uploaded_file) if uploaded_file.name.endswith('.csv') else pd.read_excel(uploaded_file)
            st.success(f"✅ Processing {df.shape[0]} rows...")
            
            ppt_content = {
                "Executive Summary": [
                    f"Dataset: {uploaded_file.name}",
                    f"Total Records: {df.shape[0]}",
                    f"Total Fields: {df.shape[1]}"
                ],
                "Column Analysis": list(df.columns),
                "Data Quality": {
                    "Complete Records": len(df.dropna()),
                    "Missing Values": df.isnull().sum().sum(),
                    "Duplicate Rows": df.duplicated().sum()
                },
                "Insights": [
                    "Data has been analyzed and loaded successfully",
                    f"Processing completed at {datetime.now().strftime('%H:%M:%S')}",
                    "Ready for presentation"
                ]
            }
            
            prs = create_ppt_from_data(ppt_content, ppt_title, (31, 119, 180))
            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            st.success("✅ Complete presentation created!")
            st.download_button(
                label="📥 Download Presentation",
                data=ppt_buffer,
                file_name=f"{ppt_title.replace(' ', '_')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except Exception as e:
            st.error(f"Error: {str(e)}")

st.sidebar.markdown("---")
st.sidebar.info("✨ **Features**:\n" +
                "- Upload & analyze Excel/CSV files\n" +
                "- Generate professional PowerPoint decks\n" +
                "- Download processed files\n" +
                "- All processing done locally")
